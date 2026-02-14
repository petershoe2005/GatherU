import os
import time
import subprocess
from pptx import Presentation
from pptx.util import Inches, Emu

# Configuration
HTML_FILE = os.path.abspath("pitchdeck.html")
CHROME_PATH = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"
OUTPUT_PPTX = "GatherU_PitchDeck.pptx"
WIDTH_PX = 1920
HEIGHT_PX = 1080

# 16:9 aspect ratio standard for PPT
PPT_WIDTH = Emu(12192000)  # 13.333 inches * 914400
PPT_HEIGHT = Emu(6858000)   # 7.5 inches * 914400

def capture_slide(index, output_filename):
    url = f"file://{HTML_FILE}?slide={index}"
    cmd = [
        CHROME_PATH,
        "--headless",
        "--hide-scrollbars",
        f"--window-size={WIDTH_PX},{HEIGHT_PX}",
        f"--screenshot={output_filename}",
        url
    ]
    subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    return output_filename

def create_pptx():
    prs = Presentation()
    prs.slide_width = PPT_WIDTH
    prs.slide_height = PPT_HEIGHT
    
    # Use blank layout (index 6 typically, but let's find it)
    blank_layout = prs.slide_layouts[6] 

    print("Generating screenshots...")
    temp_files = []
    
    try:
        # 12 slides total (indices 0 to 11)
        for i in range(12):
            filename = f"slide_temp_{i}.png"
            print(f"  Capturing slide {i+1}...")
            capture_slide(i, filename)
            temp_files.append(filename)
            
            # Add slide
            slide = prs.slides.add_slide(blank_layout)
            
            # Add image full bleed
            slide.shapes.add_picture(filename, 0, 0, PPT_WIDTH, PPT_HEIGHT)
            
        print(f"Saving {OUTPUT_PPTX}...")
        prs.save(OUTPUT_PPTX)
        print("Done!")
        
    finally:
        # Cleanup
        for f in temp_files:
            if os.path.exists(f):
                os.remove(f)

if __name__ == "__main__":
    create_pptx()
